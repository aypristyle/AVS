from os.path import exists
import speech_recognition as sr
import mss
import numpy as np
import os
from PIL import Image
path, dirs, files = next(os.walk("D:/Document/3INFO/BDD/Demon/"))
monitor =2
i = len(files)

import glob
def record_volume(path,i):
    fichier=open(path[0:-3]+".txt","a")
    r = sr.Recognizer()
    with sr.Microphone(device_index = 3) as source:
        print('.')
        r.adjust_for_ambient_noise(source, duration=0.5) #
        print('...')
        audio = r.listen(source)
    print('.')
    try:
        query = r.recognize_google(audio, language = 'fr-FR')
        text = query.lower()
        fichier.write(text+"\n")
        fichier.close()
        print(f' : {text}')
        rename(path,i)
    except:
        print('Error')
        rename(path,i)
def repartition(filename):
    image_file = Image.open(filename)
    nb = image_file.convert('1')
    tab = np.array(nb.getdata())
    nt = tab.size
    n1 = np.count_nonzero(tab == tab.max())
    return n1 / nt

def rename(path,i):
    with mss.mss() as mss_instance:
        mss_instance.shot(mon=2, output=path[0:-4]+"screen"+str(i)+"bis.png")
    if repartition(path+"screen"+str(i)+".png")-repartition(path+"screen"+str(i)+"bis.png")>-0.000001:
        print("if")
        os.remove(path+"screen"+str(i)+".png")
        os.rename(path+"screen"+str(i)+"bis.png", path)
        record_volume(path,i)
    else:
        with mss.mss() as mss_instance:
            mss_instance.shot(mon=2, output=path+"screen"+str(i+1)+".png")
        record_volume(path,i+1)

# with mss.mss() as mss_instance:
#     mss_instance.shot(mon=2, output=path)
# record_volume(path,i)
from pptx import Presentation
from pptx.util import Inches
def compile():
    image=glob.glob(path+"*png")
    print(image)
    data=glob.glob(path+"*txt")
    print(data)
    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    for img_path in image:
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = Inches(0)
        pic = slide.shapes.add_picture(img_path, left, top,height=Inches(10))
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        if exists(img_path[0:-3]+".txt"):
            text_frame.text = open(img_path[0:-3]+".txt","r").read()
    prs.save('test.pptx')
def start():
    with mss.mss() as mss_instance:
        mss_instance.shot(mon=monitor, output=path)
    record_volume(path,i)
